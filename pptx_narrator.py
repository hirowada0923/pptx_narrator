# MIT License
#
# Copyright (c) 2025 Hiro Wada
#
# Permission is hereby granted, free of charge, to any person obtaining a copy
# of this software and associated documentation files (the "Software"), to deal
# in the Software without restriction, including without limitation the rights
# to use, copy, modify, merge, publish, distribute, sublicense, and/or sell
# copies of the Software, and to permit persons to whom the Software is
# furnished to do so, subject to the following conditions:
#
# The above copyright notice and this permission notice shall be included in all
# copies or substantial portions of the Software.
#
# THE SOFTWARE IS PROVIDED "AS IS", WITHOUT WARRANTY OF ANY KIND, EXPRESS OR
# IMPLIED, INCLUDING BUT NOT LIMITED TO THE WARRANTIES OF MERCHANTABILITY,
# FITNESS FOR A PARTICULAR PURPOSE AND NONINFRINGEMENT. IN NO EVENT SHALL THE
# AUTHORS OR COPYRIGHT HOLDERS BE LIABLE FOR ANY CLAIM, DAMAGES OR OTHER
# LIABILITY, WHETHER IN AN ACTION OF CONTRACT, TORT OR OTHERWISE, ARISING FROM,
# OUT OF OR IN CONNECTION WITH THE SOFTWARE OR THE USE OR OTHER DEALINGS IN THE
# SOFTWARE.

# pptx_narrator.py
# Version 1.0
#

import logging
import os
import re
import struct
import mimetypes
import argparse
from google import genai
from google.genai import types
from mutagen.wave import WAVE
from pptx import Presentation
from lxml import etree
from pptx.util import Cm

# --- Configuration ---
# General Settings
LOG_LEVEL = logging.INFO
TEXT_TRUNCATION_LIMIT = 5000

# Gemini API Settings
TTS_MODEL = "gemini-2.5-flash-preview-tts"
DEFAULT_VOICE_1 = "Autonoe"
DEFAULT_VOICE_2 = "Algieba"

# PowerPoint Settings
# Position and size of the audio icon on the slide: (left, top, width, height)
AUDIO_SHAPE_POSITION = (Cm(30.9), Cm(16.5), Cm(3.0), Cm(1.5))


# --- Helper functions ---
def convert_to_wav(audio_data: bytes, mime_type: str) -> bytes:
    """Converts raw audio data to WAV format by adding a WAV header."""
    parsed_audio_info = parse_audio_mime_type(mime_type)
    bits_per_sample = parsed_audio_info.get("bits_per_sample", 16)
    rate = parsed_audio_info.get("rate", 24000) # Default to 24000 Hz as per Gemini TTS L16
    channels = 1  # Assuming mono for TTS output

    bytes_per_sample = bits_per_sample // 8
    block_align = channels * bytes_per_sample
    avg_bytes_per_second = rate * block_align
    data_size = len(audio_data)

    header = struct.pack(
        "<4sI4s4sIHHIIHH4sI",
        b"RIFF", 36 + data_size, b"WAVE", b"fmt ",
        16, 1, channels, rate, avg_bytes_per_second,
        block_align, bits_per_sample, b"data", data_size
    )
    return header + audio_data

def parse_audio_mime_type(mime_type: str) -> dict:
    """
    Parses bits per sample and rate from an audio MIME type string.
    Enhanced to handle L16 PCM format properly.
    """
    bits_per_sample = 16
    rate = 24000
    
    logging.debug(f"Parsing MIME type: {mime_type}")
    
    # Split by semicolon to get parameters
    parts = mime_type.split(";")
    
    # Check the main type first
    main_type = parts[0].strip().lower()
    if main_type == "audio/l16":
        bits_per_sample = 16
        logging.debug("Detected L16 format, setting bits_per_sample to 16")
    
    # Parse parameters
    for param in parts[1:]:  # Skip the first part which is the main type
        param = param.strip()
        if param.lower().startswith("rate="):
            try:
                rate = int(param.split("=", 1)[1])
                logging.debug(f"Parsed rate: {rate}")
            except (ValueError, IndexError):
                logging.warning(f"Could not parse rate from: {param}")
        elif param.lower().startswith("codec="):
            codec = param.split("=", 1)[1].lower()
            logging.debug(f"Detected codec: {codec}")
            if codec == "pcm":
                # PCM is what we expect for L16
                pass
    
    result = {"bits_per_sample": bits_per_sample, "rate": rate}
    logging.debug(f"Parsed audio info: {result}")
    return result


# --- Core Classes ---
class TextToSpeech:
    def __init__(self, model_name: str, voice_name1: str, voice_name2: str, speed: float = 1.0):
        api_key = os.environ.get("GEMINI_API_KEY")
        if not api_key:
            logging.error("API key not found. Please set the GEMINI_API_KEY environment variable.")
            raise ValueError("API key not found.")
        self.__client = genai.Client(api_key=api_key)
        self.__model = model_name
        self.__voice_name1 = voice_name1
        self.__voice_name2 = voice_name2
        self.__speaking_rate = speed

    def generate_speech_data(self, text: str) -> bytes | None:
        """
        Generates speech data from text, converts it to WAV, and returns the bytes.
        Returns None if generation or conversion fails.
        """
        if len(text) > TEXT_TRUNCATION_LIMIT:
            logging.warning(f'Speaker note is truncated > {TEXT_TRUNCATION_LIMIT} characters.')
            text = text[:TEXT_TRUNCATION_LIMIT]

        contents_parts = self._convert_text_to_gemini_content_parts(text)

        speech_config = types.SpeechConfig(
            multi_speaker_voice_config=types.MultiSpeakerVoiceConfig(
                speaker_voice_configs=[
                    types.SpeakerVoiceConfig(speaker="Speaker 1", voice_config=types.VoiceConfig(prebuilt_voice_config=types.PrebuiltVoiceConfig(voice_name=self.__voice_name1))),
                    types.SpeakerVoiceConfig(speaker="Speaker 2", voice_config=types.VoiceConfig(prebuilt_voice_config=types.PrebuiltVoiceConfig(voice_name=self.__voice_name2))),
                ]
            ),
        )
        generate_content_config = types.GenerateContentConfig(temperature=1, response_modalities=["audio"], speech_config=speech_config)

        audio_data_buffer = b''
        received_mime_type = None
        try:
            logging.info(f"Requesting TTS for: {text[:70]}...")
            for chunk in self.__client.models.generate_content_stream(model=self.__model, contents=[types.Content(role="user", parts=contents_parts)], config=generate_content_config):
                if chunk.candidates and chunk.candidates[0].content and chunk.candidates[0].content.parts and chunk.candidates[0].content.parts[0].inline_data and chunk.candidates[0].content.parts[0].inline_data.data:
                    audio_data_buffer += chunk.candidates[0].content.parts[0].inline_data.data
                    if received_mime_type is None:
                        received_mime_type = chunk.candidates[0].content.parts[0].inline_data.mime_type
                elif chunk.text:
                    logging.debug(f"Received text in stream: {chunk.text}")

            if not audio_data_buffer:
                logging.warning('No audio content was generated.')
                return None

            logging.info(f"Received audio data: {len(audio_data_buffer)} bytes, MIME type: {received_mime_type}")

            # Check if this is a format we can convert to WAV
            if received_mime_type:
                mime_lower = received_mime_type.lower()
                if ("audio/l16" in mime_lower) or ("audio/pcm" in mime_lower) or ("pcm" in mime_lower):
                    logging.info(f"Converting raw PCM audio ({received_mime_type}) to WAV format.")
                    return convert_to_wav(audio_data_buffer, received_mime_type)
                else:
                    logging.warning(f"Received unexpected audio format '{received_mime_type}'. Attempting conversion anyway.")
                    # Try to convert anyway - it might work
                    return convert_to_wav(audio_data_buffer, received_mime_type)
            else:
                logging.warning("No MIME type received, assuming L16 PCM format.")
                return convert_to_wav(audio_data_buffer, "audio/L16;rate=24000")
            
        except Exception as e:
            logging.error(f"Error during TTS generation: {e}", exc_info=True)
            return None

    def _convert_text_to_gemini_content_parts(self, text: str) -> list:
        content_parts = []
        for line in text.splitlines():
            line = line.strip()
            if not line:
                continue
            match1 = re.match(r'Speaker 1:\s*(.*)', line, re.IGNORECASE)
            match2 = re.match(r'Speaker 2:\s*(.*)', line, re.IGNORECASE)
            if match1:
                content_parts.append(types.Part.from_text(text=f'Speaker 1: {match1.group(1).strip()}'))
            elif match2:
                content_parts.append(types.Part.from_text(text=f'Speaker 2: {match2.group(1).strip()}'))
            else:
                content_parts.append(types.Part.from_text(text=f'Speaker 1: {line}'))
        if not content_parts and text.strip():
            content_parts.append(types.Part.from_text(text=f'Speaker 1: {text.strip()}'))
        return content_parts

class Powerpoint:
    def __init__(self, pptx_path: str, tts_engine=None):
        self.__fname = pptx_path
        self.__name  = os.path.splitext(pptx_path)[0]
        self.__pp    = Presentation(pptx_path)
        self.__tts   = tts_engine

    def VoiceAnnotatePP(self, out_pptx_path: str):
        if not self.__tts:
            logging.error("TTS engine not provided. Aborting.")
            return

        total_slides = len(self.__pp.slides)
        for idx, s in enumerate(self.__pp.slides):
            logging.info(f'Processing slide {idx+1}/{total_slides}...')
            if not s.has_notes_slide or not s.notes_slide.notes_text_frame.text.strip():
                logging.info('No speaker notes found.')
                continue

            text = s.notes_slide.notes_text_frame.text
            wav_data = self.__tts.generate_speech_data(text)
            if not wav_data:
                logging.warning(f"Skipping audio for slide {idx+1} due to generation failure.")
                continue

            outf_path = f'{self.__name}-{str(idx+1):>03}.wav'
            try:
                with open(outf_path, 'wb') as out:
                    out.write(wav_data)
                logging.info(f'Created "{outf_path}"')
                try:
                    audio = WAVE(outf_path)
                    logging.info(f'  -> Length: {audio.info.length:.2f}s')
                except Exception as mutagen_e:
                    logging.warning(f'Could not get audio length: {mutagen_e}')
            except IOError as e:
                logging.error(f"Error writing audio file '{outf_path}': {e}")
                continue

            try:
                movie = s.shapes.add_movie(outf_path, *AUDIO_SHAPE_POSITION, poster_frame_image=None, mime_type='audio/wav')
                tree = movie._element.getparent().getparent().getnext().getnext().getnext()
                if tree is not None:
                    timing_elements = [el for el in tree.iterdescendants() if etree.QName(el).localname == 'cond']
                    if timing_elements:
                        timing_elements[0].set('delay', '0')
                else: # Fallback for different structures
                    tree = movie._element.getparent().getparent().getnext().getnext()
                    if tree is not None:
                        timing_elements = [el for el in tree.iterdescendants() if etree.QName(el).localname == 'cond']
                        if timing_elements:
                            timing_elements[0].set('delay', '0')
            except Exception as e:
                logging.error(f"Error adding audio to slide {idx+1}: {e}", exc_info=True)
        
        logging.info(f"Saving annotated presentation to {out_pptx_path}...")
        try:
            self.__pp.save(out_pptx_path)
            logging.info("Successfully saved.")
        except Exception as e:
            logging.error(f"Error saving PowerPoint file: {e}", exc_info=True)


# --- Main Execution ---
if __name__ == '__main__':
    logging.basicConfig(level=LOG_LEVEL, format='%(asctime)s - %(levelname)s - %(message)s')

    parser = argparse.ArgumentParser(description='Voice annotate PowerPoint file with multi-speaker support using Gemini API TTS')
    parser.add_argument('source_pptx', type=str, help='Input PowerPoint file')
    parser.add_argument('output_pptx', type=str, help='Output PowerPoint file')
    parser.add_argument('--name1', nargs='?', default=DEFAULT_VOICE_1, help=f'Voice name for Speaker 1 (default: {DEFAULT_VOICE_1})')
    parser.add_argument('--name2', nargs='?', default=DEFAULT_VOICE_2, help=f'Voice name for Speaker 2 (default: {DEFAULT_VOICE_2})')
    parser.add_argument('--speed',  nargs='?', default=1.0, help='Speaking speed for all speakers (default: 1.0)', type=float)
    args = parser.parse_args()

    logging.info(f"Source PPTX: {args.source_pptx}")
    logging.info(f"Output PPTX: {args.output_pptx}")
    logging.info(f"Speaker 1 Voice: {args.name1}")
    logging.info(f"Speaker 2 Voice: {args.name2}")
    logging.info(f"Overall Speed: {args.speed}")

    try:
        tts_engine = TextToSpeech(
            model_name=TTS_MODEL,
            voice_name1=args.name1,
            voice_name2=args.name2,
            speed=args.speed
        )
        ppt_processor = Powerpoint(args.source_pptx, tts_engine)
        ppt_processor.VoiceAnnotatePP(args.output_pptx)
        logging.info('Processing Done.')
    except Exception as e:
        logging.critical(f"A critical error occurred in the main execution block: {e}", exc_info=True)
